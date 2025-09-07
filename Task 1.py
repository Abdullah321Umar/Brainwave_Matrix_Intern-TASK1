import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

# CONFIGURATION / PATHS
DATA_PATH = r"C:/Users/Abdullah Umer/Desktop/Brainwave Matrix Solution Internship/Task 1/Walmart_DataSet.csv"
OUTPUT_DIR = r"C:/Users/Abdullah Umer/Desktop/Brainwave Matrix Solution Internship/Task 1/outputs"
PPTX_PATH = os.path.join(OUTPUT_DIR, "Walmart_Sales_Report.pptx")

os.makedirs(OUTPUT_DIR, exist_ok=True)
print("🎯 Saving outputs to:", OUTPUT_DIR)

# HELPER FUNCTION
def save_fig(fig, name, dpi=150):
    """Save matplotlib figure into outputs folder."""
    path = os.path.join(OUTPUT_DIR, name)
    fig.savefig(path, bbox_inches="tight", dpi=dpi)
    plt.close(fig)
    return path

# 1) DATA COLLECTION & LOADING
print("📌 Loading data from:", DATA_PATH)
df = pd.read_csv(DATA_PATH)
print("📌 Initial shape:", df.shape)
print("✅ Columns:", df.columns.tolist())
print(df.head())

# 2) DATA CLEANING
date_col = "Date"

# Convert Date column
df[date_col] = pd.to_datetime(df[date_col], dayfirst=True, errors="coerce")

# Drop duplicates
df = df.drop_duplicates()

# Standardize column names
df.columns = [c.strip() for c in df.columns]

# Convert Weekly_Sales if needed
if df["Weekly_Sales"].dtype == "O":
    df["Weekly_Sales"] = pd.to_numeric(df["Weekly_Sales"].str.replace(",", ""), errors="coerce")

# Create new features
df["Year"] = df[date_col].dt.year
df["Month"] = df[date_col].dt.month
df["Week"] = df[date_col].dt.isocalendar().week
df["DayOfWeek"] = df[date_col].dt.dayofweek

# Handle missing values
threshold = 0.05 * len(df)
for col in df.columns:
    if df[col].isnull().sum() > 0:
        if df[col].dtype in [np.float64, np.int64]:
            df[col] = df[col].fillna(df[col].median())
        else:
            if df[col].isnull().sum() < threshold:
                df = df.dropna(subset=[col])
            else:
                df[col] = df[col].fillna("Unknown")

print("✅ Shape after cleaning:", df.shape)

# 3) SUMMARY STATS
total_sales = df["Weekly_Sales"].sum()
mean_sales = df["Weekly_Sales"].mean()
median_sales = df["Weekly_Sales"].median()

sales_by_date = df.groupby(date_col)["Weekly_Sales"].sum().sort_index()
sales_by_store = df.groupby("Store")["Weekly_Sales"].sum().sort_values(ascending=False)
avg_sales_by_store = df.groupby("Store")["Weekly_Sales"].mean().sort_values(ascending=False)
sales_by_year = df.groupby("Year")["Weekly_Sales"].sum().sort_index()
sales_by_month = df.groupby("Month")["Weekly_Sales"].sum().sort_index()
sales_holiday_flag = df.groupby("Holiday_Flag")["Weekly_Sales"].mean()

print(f"📌 Total sales: {total_sales:,.2f}")
print(f"📌 Mean weekly sales: {mean_sales:,.2f}")
print(f"📌 Top store by total sales: {sales_by_store.index[0]} -> {sales_by_store.iloc[0]:,.2f}")

# 4) VISUALIZATIONS
sns.set(style="whitegrid", rc={"figure.figsize": (12, 6)})

# 4.1 Line plot - Weekly Sales Trend
fig, ax = plt.subplots()
sales_by_date.plot(ax=ax)
ax.set_title("Weekly Sales Trend (Total across all stores)")
ax.set_xlabel("Date")
ax.set_ylabel("Total Weekly Sales")
line_path = save_fig(fig, "01_sales_trend_line.png")

# 4.2 Bar plot - Top 10 Stores
fig, ax = plt.subplots()
top10_avg = avg_sales_by_store.head(10)
top10_avg.plot(kind="bar", ax=ax)
ax.set_title("Top 10 Stores by Average Weekly Sales")
ax.set_ylabel("Average Weekly Sales")
bar_path = save_fig(fig, "02_top10_avg_store_bar.png")

# 4.3 Pie chart - Holiday vs Non-Holiday Sales
fig, ax = plt.subplots(figsize=(6, 6))
labels = ["Non-Holiday", "Holiday"]
vals = [sales_holiday_flag.get(0, 0), sales_holiday_flag.get(1, 0)]
ax.pie(vals, labels=labels, autopct="%1.1f%%", startangle=90)
ax.set_title("Average Weekly Sales: Holiday vs Non-Holiday")
pie_path = save_fig(fig, "03_holiday_pie.png")

# 4.4 Heatmap - Correlation
fig, ax = plt.subplots(figsize=(10, 8))
numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
corr = df[numeric_cols].corr()
sns.heatmap(corr, annot=True, fmt=".2f", cmap="coolwarm", ax=ax)
ax.set_title("Correlation Heatmap (Numeric Features)")
heatmap_path = save_fig(fig, "04_correlation_heatmap.png")

# 4.5 Histogram - Weekly Sales
fig, ax = plt.subplots()
sns.histplot(df["Weekly_Sales"], bins=50, kde=True, ax=ax)
ax.set_title("Distribution of Weekly Sales")
hist_path = save_fig(fig, "05_sales_histogram.png")

# 4.6 Boxplot - Sales by Year
fig, ax = plt.subplots(figsize=(12, 6))
sns.boxplot(x="Year", y="Weekly_Sales", data=df, ax=ax)
ax.set_title("Sales Distribution by Year")
box_path = save_fig(fig, "06_boxplot_year.png")

# 4.7 Scatter - Temperature vs Sales
if "Temperature" in df.columns:
    fig, ax = plt.subplots()
    sns.scatterplot(x="Temperature", y="Weekly_Sales", hue="Holiday_Flag", data=df, ax=ax)
    ax.set_title("Temperature vs Weekly Sales")
    scatter_temp_path = save_fig(fig, "07_scatter_temperature_sales.png")
else:
    scatter_temp_path = None

# 4.8 Scatter - Fuel Price vs Sales
if "Fuel_Price" in df.columns:
    fig, ax = plt.subplots()
    sns.scatterplot(x="Fuel_Price", y="Weekly_Sales", data=df, ax=ax)
    ax.set_title("Fuel Price vs Weekly Sales")
    scatter_fuel_path = save_fig(fig, "08_scatter_fuel_sales.png")
else:
    scatter_fuel_path = None

# 4.9 Area Chart - Monthly Sales
monthly = df.set_index(date_col).resample("ME")["Weekly_Sales"].sum()
fig, ax = plt.subplots()
monthly.plot(kind="area", ax=ax)
ax.set_title("Monthly Total Sales (Area)")
area_path = save_fig(fig, "09_monthly_area.png")

# 5) INSIGHTS
insights = []
insights.append(f"Total sales (dataset): {total_sales:,.2f}")
insights.append(f"Average weekly sales: {mean_sales:,.2f}")
insights.append(f"Median weekly sales: {median_sales:,.2f}")
insights.append(f"Top store by total sales: Store {sales_by_store.index[0]} with {sales_by_store.iloc[0]:,.2f}")
insights.append(f"Top store by average weekly sales: Store {avg_sales_by_store.index[0]} with {avg_sales_by_store.iloc[0]:,.2f}")
insights.append(f"Holiday vs Non-Holiday average weekly sales: Non-Holiday {sales_holiday_flag.get(0,0):,.2f}, Holiday {sales_holiday_flag.get(1,0):,.2f}")
if not monthly.empty:
    highest_month = monthly.idxmax()
    insights.append(f"Highest monthly sales observed in: {highest_month.strftime('%Y-%m')} with {monthly.max():,.2f}")

print("\n🌟 ======== Key Insights ======== 🌟")
for i, s in enumerate(insights, 1):
    print(f"{i}. {s}")

# 6) POWERPOINT REPORT
print("\n📌 Generating PowerPoint report...")

prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Walmart Sales Analysis"
subtitle.text = "Automated Report Generated by Python\n(Charts & Insights)"

# Function to add image slide
def add_image_slide(prs, title_text, image_path, caption=None):
    slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text

    left = Inches(1)
    top = Inches(1.6)
    height = Inches(4.5)

    try:
        slide.shapes.add_picture(image_path, left, top, height=height)
    except Exception as e:
        txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = f"Could not insert image: {e}"

    if caption:
        tx = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
        tf = tx.text_frame
        tf.text = caption

# Add charts to PowerPoint
add_image_slide(prs, "Weekly Sales Trend", line_path)
add_image_slide(prs, "Top 10 Stores", bar_path)
add_image_slide(prs, "Holiday vs Non-Holiday Sales", pie_path)
add_image_slide(prs, "Correlation Heatmap", heatmap_path)
add_image_slide(prs, "Sales Distribution", hist_path)
add_image_slide(prs, "Sales by Year", box_path)
add_image_slide(prs, "Monthly Sales", area_path)
if scatter_temp_path: add_image_slide(prs, "Temperature vs Sales", scatter_temp_path)
if scatter_fuel_path: add_image_slide(prs, "Fuel Price vs Sales", scatter_fuel_path)

# Insights slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Key Insights & Recommendations"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame

for i, s in enumerate(insights, 1):
    p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
    p.text = f"{i}. {s}"
    p.level = 0

# Save PPTX
prs.save(PPTX_PATH)
print("✅ PowerPoint saved to:", PPTX_PATH)

# Save cleaned dataset
clean_path = os.path.join(OUTPUT_DIR, "Walmart_cleaned.csv")
df.to_csv(clean_path, index=False)
print("✅ Cleaned dataset saved to:", clean_path)

print("\nAll done. Outputs in:", OUTPUT_DIR)






